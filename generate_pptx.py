#!/usr/bin/env python3
"""Generate All Shanti UNICEF Pitch PowerPoint from slide components."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# Color palette
NAVY = RGBColor(0x0A, 0x44, 0xA1)
UNICEF_BLUE = RGBColor(0x00, 0x9E, 0xDB)
PINK = RGBColor(0xE8, 0x6E, 0xA6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1a, 0x20, 0x2c)
GRAY_50 = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_100 = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_400 = RGBColor(0x9C, 0xA3, 0xAF)
GRAY_500 = RGBColor(0x6B, 0x72, 0x80)
GRAY_600 = RGBColor(0x4B, 0x55, 0x63)
BLUE_50 = RGBColor(0xEF, 0xF6, 0xFF)
BLUE_100 = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_200 = RGBColor(0xBF, 0xDB, 0xFE)
RED_50 = RGBColor(0xFF, 0xF1, 0xF2)
RED_500 = RGBColor(0xEF, 0x44, 0x44)
GREEN_50 = RGBColor(0xF0, 0xFF, 0xF4)
GREEN_600 = RGBColor(0x16, 0xA3, 0x4A)
AMBER_50 = RGBColor(0xFF, 0xFB, 0xEB)
AMBER_600 = RGBColor(0xD9, 0x77, 0x06)
PINK_50 = RGBColor(0xFF, 0xF0, 0xF7)

# Slide dimensions: 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_background_fill(slide, color: RGBColor):
    """Fill the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_gradient_background(slide, color1: RGBColor, color2: RGBColor):
    """Add a gradient background using XML manipulation."""
    # Use a rectangle covering the whole slide as background
    from pptx.util import Pt
    left = top = 0
    width = SLIDE_W
    height = SLIDE_H
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.AUTO_SHAPE if False else 1,  # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
        left, top, width, height
    )
    shape.fill.gradient()
    shape.fill.gradient_angle = 135
    shape.fill.gradient_stops[0].position = 0.0
    shape.fill.gradient_stops[0].color.rgb = color1
    shape.fill.gradient_stops[1].position = 1.0
    shape.fill.gradient_stops[1].color.rgb = color2
    shape.line.fill.background()
    # Send to back
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def add_rect(slide, left, top, width, height, fill_color: RGBColor, line_color=None):
    """Add a filled rectangle."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=12, bold=False,
                color=DARK, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Add a text box with specified properties."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox


def add_header(slide, label, title, bg_color1=NAVY, bg_color2=None, subtitle=None):
    """Add a slide header with label + title."""
    header_h = Inches(1.3)

    if bg_color2:
        # Gradient header
        rect = add_rect(slide, 0, 0, SLIDE_W, header_h, bg_color1)
        rect.fill.gradient()
        rect.fill.gradient_angle = 0
        rect.fill.gradient_stops[0].position = 0.0
        rect.fill.gradient_stops[0].color.rgb = bg_color1
        rect.fill.gradient_stops[1].position = 1.0
        rect.fill.gradient_stops[1].color.rgb = bg_color2
        rect.line.fill.background()
    else:
        add_rect(slide, 0, 0, SLIDE_W, header_h, bg_color1)

    # Category label
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.3),
                label, font_size=8, bold=True, color=BLUE_200)

    # Title
    add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
                title, font_size=22, bold=True, color=WHITE)

    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.3),
                    subtitle, font_size=10, color=BLUE_200)

    return header_h


def add_source(slide, text):
    """Add source citation at the bottom."""
    add_textbox(slide, Inches(0.8), Inches(7.1), Inches(12), Inches(0.3),
                text, font_size=7, color=GRAY_400)


# ============================================================
# Slide 01 — Cover
# ============================================================
def slide01_cover(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # Gradient background
    add_gradient_background(slide, NAVY, UNICEF_BLUE)

    # Pink circle decoration top-right (subtle)
    circ = add_rect(slide, Inches(10), Inches(-1), Inches(4), Inches(4), PINK)
    # Make it a circle shape (rounded)

    # Center content
    # Logo circle with "AS"
    logo_x = Inches(5.9)
    logo_y = Inches(0.9)
    logo_size = Inches(0.9)
    logo = add_rect(slide, logo_x, logo_y, logo_size, logo_size, PINK)

    # AS text on logo
    add_textbox(slide, logo_x, logo_y, logo_size, logo_size,
                "AS", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Main title
    add_textbox(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(1.0),
                "All Shanti", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(0.7),
                "Proposta de Investigação Clínica no Contexto de Tratamento de Suporte Contra a Malária",
                font_size=12, color=BLUE_200, align=PP_ALIGN.CENTER)

    # Pink line divider
    line = add_rect(slide, Inches(5.9), Inches(3.85), Inches(1.5), Inches(0.06), PINK)

    # Company name
    add_textbox(slide, Inches(1), Inches(4.1), Inches(11.3), Inches(0.4),
                "Águas de São Silvestre, SA", font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Presentation line
    add_textbox(slide, Inches(1), Inches(4.5), Inches(11.3), Inches(0.4),
                "Apresentação à UNICEF • 2026", font_size=9, color=BLUE_200, align=PP_ALIGN.CENTER)

    # Confidential badge
    add_textbox(slide, Inches(1), Inches(7.0), Inches(11.3), Inches(0.35),
                "CONFIDENCIAL", font_size=8, color=BLUE_200, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 02 — Crisis
# ============================================================
def slide02_crisis(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "CONTEXTO", "A Crise de Malária em África", bg_color1=NAVY)

    stats = [
        ('282M', 'casos globais em 2024', UNICEF_BLUE),
        ('610K', 'mortes registadas em 2024', PINK),
        ('95%', 'dos casos mundiais ocorrem em África', NAVY),
        ('75%', 'das vítimas mortais são crianças < 5 anos', PINK),
    ]

    card_w = Inches(5.5)
    card_h = Inches(2.1)
    positions = [
        (Inches(0.6), Inches(1.5)),
        (Inches(6.9), Inches(1.5)),
        (Inches(0.6), Inches(3.9)),
        (Inches(6.9), Inches(3.9)),
    ]

    for (stat_val, stat_label, stat_color), (x, y) in zip(stats, positions):
        add_rect(slide, x, y, card_w, card_h, GRAY_50)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.2), card_w - Inches(0.4), Inches(0.9),
                    stat_val, font_size=42, bold=True, color=stat_color)
        add_textbox(slide, x + Inches(0.3), y + Inches(1.1), card_w - Inches(0.4), Inches(0.6),
                    stat_label, font_size=11, color=DARK)

    # Context paragraph
    context = ("A malária é causada por parasitas Plasmodium, transmitidos pela picada de mosquitos Anopheles fêmea. "
               "As duas espécies principais são P. falciparum — a mais letal, dominante em África — e P. vivax, "
               "a mais comum fora de África subsariana. A distribuição geográfica concentra-se em África (~94% dos casos), "
               "seguida do Sudeste Asiático (~2%) e das regiões do Mediterrâneo Oriental e Américas.")
    add_textbox(slide, Inches(0.6), Inches(6.15), Inches(12.1), Inches(0.65),
                context, font_size=8, color=GRAY_500)

    add_source(slide, "Fonte: WHO World Malaria Report 2025")


# ============================================================
# Slide 03 — Countries
# ============================================================
def slide03_countries(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "EPIDEMIOLOGIA", "Países Mais Afetados",
                           bg_color1=UNICEF_BLUE,
                           subtitle="% das mortes na região africana — WHO 2025")

    countries = [
        ('Nigéria', 31.9, UNICEF_BLUE),
        ('RDCongo', 11.7, PINK),
        ('Niger', 6.1, NAVY),
        ('Uganda, Moçambique e outros', 50.3, GRAY_400),
    ]

    bar_area_left = Inches(2.5)
    bar_area_width = Inches(10.0)
    bar_h = Inches(0.4)

    for i, (name, pct, color) in enumerate(countries):
        y = Inches(1.5) + i * Inches(0.75)
        # Country label
        add_textbox(slide, Inches(0.3), y, Inches(2.0), bar_h + Inches(0.1),
                    name, font_size=10, bold=True, color=DARK, align=PP_ALIGN.RIGHT)
        # Background bar
        add_rect(slide, bar_area_left, y + Inches(0.03), bar_area_width, bar_h, GRAY_100)
        # Filled bar (scale: pct * 2.5 / 100 of bar area)
        filled_w = max(Inches(0.5), bar_area_width * (pct * 2.5 / 100))
        filled_rect = add_rect(slide, bar_area_left, y + Inches(0.03), filled_w, bar_h, color)
        # Percentage label on bar
        add_textbox(slide, bar_area_left + filled_w - Inches(0.7), y + Inches(0.03),
                    Inches(0.65), bar_h, f"{pct}%", font_size=9, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    # Blue info box
    blue_y = Inches(4.65)
    add_rect(slide, Inches(0.6), blue_y, Inches(12.1), Inches(0.55), BLUE_50)
    add_textbox(slide, Inches(0.9), blue_y + Inches(0.08), Inches(11.6), Inches(0.4),
                "3 países — Nigéria, RDCongo e Niger — representam ~50% de todas as mortes por malária em África.",
                font_size=10, bold=False, color=NAVY)

    # Pink info box
    pink_y = Inches(5.3)
    add_rect(slide, Inches(0.6), pink_y, Inches(12.1), Inches(0.85), PINK_50)
    add_textbox(slide, Inches(0.9), pink_y + Inches(0.05), Inches(11.6), Inches(0.3),
                "78% das mortes por malária são crianças com menos de 5 anos — a cada minuto, uma criança morre de malária.",
                font_size=10, bold=True, color=PINK)
    add_textbox(slide, Inches(0.9), pink_y + Inches(0.38), Inches(11.6), Inches(0.38),
                "As vacinas RTS,S e R21 estão agora em fase de rollout em mais de 20 países africanos, com resultados significativos na redução da mortalidade grave.",
                font_size=9, color=GRAY_600)

    add_source(slide, "Fonte: WHO World Malaria Report 2025")


# ============================================================
# Slide 04 — Economic Impact
# ============================================================
def slide04_economic(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "IMPACTO", "Impacto Económico Devastador",
                           bg_color1=NAVY, bg_color2=RGBColor(0x1a, 0x5b, 0xd4))

    # Left card — Cost
    red_bg = RGBColor(0xFF, 0xF1, 0xF2)
    card_h = Inches(2.2)
    card_y = Inches(1.5)

    add_rect(slide, Inches(0.6), card_y, Inches(5.8), card_h, red_bg)
    add_textbox(slide, Inches(0.85), card_y + Inches(0.12), Inches(5.3), Inches(0.25),
                "CUSTO ANUAL DA MALÁRIA", font_size=7, bold=True, color=RED_500)
    add_textbox(slide, Inches(0.85), card_y + Inches(0.35), Inches(5.3), Inches(0.55),
                "US$ 12 mil M", font_size=28, bold=True, color=RED_500)
    add_textbox(slide, Inches(0.85), card_y + Inches(0.9), Inches(5.3), Inches(0.3),
                "em perdas directas por ano em África", font_size=9, color=GRAY_600)
    add_textbox(slide, Inches(0.85), card_y + Inches(1.25), Inches(5.3), Inches(0.7),
                "• Perdas de produtividade laboral\n• Custos de saúde e investimento perdido",
                font_size=9, color=GRAY_600)

    # Right card — Gap
    green_bg = RGBColor(0xF0, 0xFF, 0xF4)
    add_rect(slide, Inches(6.9), card_y, Inches(5.8), card_h, green_bg)
    add_textbox(slide, Inches(7.15), card_y + Inches(0.12), Inches(5.3), Inches(0.25),
                "GAP DE FINANCIAMENTO (2023)", font_size=7, bold=True, color=GREEN_600)
    add_textbox(slide, Inches(7.15), card_y + Inches(0.35), Inches(5.3), Inches(0.55),
                "US$ 4,3 mil M", font_size=28, bold=True, color=GREEN_600)
    add_textbox(slide, Inches(7.15), card_y + Inches(0.9), Inches(5.3), Inches(0.3),
                "em falta para atingir metas WHO 2030", font_size=9, color=GRAY_600)
    add_textbox(slide, Inches(7.15), card_y + Inches(1.25), Inches(5.3), Inches(0.7),
                "• Fundos disponíveis: US$ 4 mil M em 2023\n• Necessidade estimada: US$ 8,3 mil M/ano",
                font_size=9, color=GRAY_600)

    # Social impact row
    amber_bg = AMBER_50
    soc_y = Inches(3.85)
    soc_h = Inches(1.8)
    add_rect(slide, Inches(0.6), soc_y, Inches(12.1), soc_h, amber_bg)
    add_textbox(slide, Inches(0.9), soc_y + Inches(0.1), Inches(11), Inches(0.25),
                "IMPACTO SOCIAL", font_size=7, bold=True, color=AMBER_600)

    col_w = Inches(3.7)
    for i, (val, text) in enumerate([
        ("~608 000", "mortes anuais — concentradas em populações com acesso limitado a cuidados de saúde, sobretudo crianças e grávidas"),
        ("17,8M", "mulheres grávidas assistidas em 2025 através de intervenções de prevenção e tratamento"),
        ("2050", "meta de erradicação global — requer inovação urgente em redes mosquiteiras de nova geração e novos fármacos face à resistência crescente e ao impacto de conflitos"),
    ]):
        x = Inches(0.85) + i * col_w
        add_textbox(slide, x, soc_y + Inches(0.38), col_w - Inches(0.2), Inches(0.45),
                    val, font_size=18, bold=True, color=AMBER_600)
        add_textbox(slide, x, soc_y + Inches(0.85), col_w - Inches(0.2), Inches(0.75),
                    text, font_size=8, color=GRAY_600)

    # Bottom banner
    add_rect(slide, Inches(0.6), Inches(5.75), Inches(12.1), Inches(0.55), NAVY)
    add_textbox(slide, Inches(0.9), Inches(5.83), Inches(11.6), Inches(0.4),
                "Investigação em novas abordagens terapêuticas é  URGENTE E SUBFINANCIADA  face à resistência crescente",
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_source(slide, "Fonte: RBM Partnership to End Malaria; WHO World Malaria Report 2024; The Lancet Commission; Global Fund Results Report 2025")


# ============================================================
# Slide 05 — Treatment
# ============================================================
def slide05_treatment(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "PROBLEMA", "Custos e Limitações dos Tratamentos Atuais",
                           bg_color1=UNICEF_BLUE)

    items = [
        ('Artemisina (standard)', 'US$ 1–3', 'por tratamento completo', UNICEF_BLUE),
        ('ACT combinadas', 'US$ 0,25–1,50', 'por dose adulto no sector público', NAVY),
        ('Acesso limitado', '43%', 'da pop. subsaariana sem proteção adequada', PINK),
        ('Resistência crescente', '8 países', 'com resistência parcial confirmada ou suspeita em África (WHO 2025)', RED_500),
    ]

    positions = [
        (Inches(0.6), Inches(1.5)),
        (Inches(6.9), Inches(1.5)),
        (Inches(0.6), Inches(4.2)),
        (Inches(6.9), Inches(4.2)),
    ]

    card_w = Inches(5.8)
    card_h = Inches(2.3)

    for (title, value, desc, color), (x, y) in zip(items, positions):
        add_rect(slide, x, y, Inches(0.08), card_h, color)
        add_rect(slide, x + Inches(0.08), y, card_w - Inches(0.08), card_h, GRAY_50)
        add_textbox(slide, x + Inches(0.35), y + Inches(0.2), card_w - Inches(0.5), Inches(0.25),
                    title.upper(), font_size=8, bold=True, color=GRAY_500)
        add_textbox(slide, x + Inches(0.35), y + Inches(0.5), card_w - Inches(0.5), Inches(0.7),
                    value, font_size=26, bold=True, color=DARK)
        add_textbox(slide, x + Inches(0.35), y + Inches(1.25), card_w - Inches(0.5), Inches(0.75),
                    desc, font_size=9, color=GRAY_600)

    add_source(slide, "Fonte: WHO Malaria Treatment Guidelines 2025; WHO World Malaria Report 2025")


# ============================================================
# Slide 06 — Solution (full-color)
# ============================================================
def slide06_solution(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Gradient background pink to navy
    add_gradient_background(slide, PINK, NAVY)

    # Main content
    add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10.3), Inches(1.1),
                "All Shanti", font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(1.85), Inches(10.3), Inches(0.45),
                "Uma Candidatura à Investigação Clínica", font_size=16, color=WHITE, align=PP_ALIGN.CENTER)

    # Pink line
    add_rect(slide, Inches(5.9), Inches(2.4), Inches(1.5), Inches(0.07), RGBColor(0xFF, 0xFF, 0xFF))

    desc = ("Suplemento alimentar desenvolvido pela Águas de São Silvestre, SA, com base em evidência histórica "
            "da Cinchona calisaya — proposto para investigação clínica formal em parceria internacional.")
    add_textbox(slide, Inches(1.5), Inches(2.65), Inches(10.3), Inches(0.8),
                desc, font_size=13, color=WHITE, align=PP_ALIGN.CENTER)

    # Three info boxes
    info_items = [
        ('BASE', 'Água Mineral Natural'),
        ('ORIGEM', 'Plantas medicinais naturais'),
        ('LOGÍSTICA', 'Sem cadeia de frio'),
    ]
    box_w = Inches(3.6)
    box_h = Inches(0.85)
    box_y = Inches(3.65)

    for i, (label, value) in enumerate(info_items):
        bx = Inches(0.7) + i * Inches(4.1)
        box2 = add_rect(slide, bx, box_y, box_w, box_h, RGBColor(0x1A, 0x44, 0x90))
        box2.line.fill.background()

        add_textbox(slide, bx + Inches(0.1), box_y + Inches(0.05), box_w - Inches(0.2), Inches(0.3),
                    label, font_size=7, color=RGBColor(0xCC, 0xCC, 0xFF), align=PP_ALIGN.CENTER)
        add_textbox(slide, bx + Inches(0.1), box_y + Inches(0.38), box_w - Inches(0.2), Inches(0.4),
                    value, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Bottom line
    add_rect(slide, Inches(5.9), Inches(4.65), Inches(1.5), Inches(0.07), WHITE)


# ============================================================
# Slide 07 — Advantages
# ============================================================
def slide07_advantages(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "PRODUTO", "Vantagens do All Shanti",
                           bg_color1=PINK, bg_color2=RGBColor(0xC4, 0x54, 0x90))

    advantages = [
        ('Evidência histórica documentada',
         'A Cinchona calisaya tem uso antimalárico documentado desde o século XVII — base para protocolo de investigação'),
        ('Formulação singular',
         'Combinação única de ingredientes naturais em água mineral de São Silvestre — a investigar em ensaios controlados'),
        ('Perfil de segurança a determinar',
         'Testes de segurança e tolerância — incluindo populações pediátricas — a realizar no âmbito do programa clínico'),
        ('Potencial de escalabilidade',
         'Se validado clinicamente, produção escalável para contexto africano de alta carga'),
        ('Sem cadeia de frio',
         'Distribuição simplificada — vantagem logística relevante para regiões remotas, sujeita a validação operacional'),
        ('Revisão bibliográfica em curso',
         'Literatura sobre alcalóides quinínicos como fundamento do protocolo de ensaios fase I/II'),
    ]

    card_w = Inches(3.9)
    card_h = Inches(1.7)

    for i, (title, desc) in enumerate(advantages):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * Inches(4.28)
        y = Inches(1.45) + row * Inches(2.0)

        add_rect(slide, x, y, card_w, card_h, GRAY_50)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.3), Inches(0.35),
                    title, font_size=10, bold=True, color=DARK)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.3), Inches(0.95),
                    desc, font_size=8, color=GRAY_500)


# ============================================================
# Slide Science — Scientific Evidence
# ============================================================
def slide_science(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "EVIDÊNCIA", "Base Científica Publicada",
                           bg_color1=NAVY, bg_color2=UNICEF_BLUE,
                           subtitle="O que a literatura peer-reviewed documenta sobre os ingredientes activos")

    evidence = [
        {
            'ingredient': 'Cinchona calisaya',
            'color': UNICEF_BLUE,
            'level': 'Forte',
            'bg': BLUE_50,
            'claims': [
                'Fonte botânica da quinina — até 6,5% de alcalóides totais (Cano et al., Frontiers in Plant Science, 2017)',
                'Quinina: mecanismo de acção contra P. falciparum documentado (inibição da hemozoína)',
                'Quinina na Lista de Medicamentos Essenciais da OMS — 2ª linha para malária grave',
                'Uso antimalárico documentado desde o séc. XVII — revista JRSM (Gachelin et al., 2017)',
            ],
            'gap': 'Sem ensaios modernos do extracto completo (não da quinina isolada)',
        },
        {
            'ingredient': "Pau d'Arco (Tabebuia)",
            'color': PINK,
            'level': 'Moderada',
            'bg': PINK_50,
            'claims': [
                'Lapachol e beta-lapachona: actividade antiplasmodial in vitro documentada',
                'Mecanismo: inibição da cadeia respiratória mitocondrial do parasita (análogo ao atovaquone)',
                'Andrade-Neto et al., Phytomedicine (2004) — IC50 1,67–9,44 μM vs. P. falciparum resistente',
                'Revisão sistemática: MDPI Molecules (2020) — PMC7571111',
            ],
            'gap': 'Sem ensaios clínicos humanos; selectividade vs. células humanas a confirmar',
        },
        {
            'ingredient': 'Lycopodium clavatum',
            'color': GRAY_400,
            'level': 'Limitada',
            'bg': GRAY_50,
            'claims': [
                'Actividade antioxidante e anti-inflamatória documentada em literatura revisada',
                'Uso tradicional como tónico digestivo e adaptogénico',
            ],
            'gap': 'Sem evidência antimalárica publicada — posicionado como ingrediente de suporte',
        },
    ]

    card_h = Inches(1.55)
    for i, ev in enumerate(evidence):
        y = Inches(1.5) + i * Inches(1.7)
        add_rect(slide, Inches(0.5), y, Inches(12.3), card_h, ev['bg'])

        # Ingredient name
        add_textbox(slide, Inches(0.7), y + Inches(0.1), Inches(2.2), Inches(0.35),
                    ev['ingredient'], font_size=10, bold=True, italic=True, color=ev['color'])
        # Level badge
        level_text = f"Evidência {ev['level']}"
        add_textbox(slide, Inches(0.7), y + Inches(0.5), Inches(2.2), Inches(0.3),
                    level_text, font_size=8, bold=True, color=ev['color'])

        # Claims
        claims_text = "\n".join([f"› {c}" for c in ev['claims']])
        add_textbox(slide, Inches(3.1), y + Inches(0.1), Inches(9.5), Inches(1.1),
                    claims_text, font_size=8, color=DARK)

        # Gap
        gap_text = f"Lacuna: {ev['gap']}"
        add_textbox(slide, Inches(3.1), y + Inches(1.2), Inches(9.5), Inches(0.28),
                    gap_text, font_size=7.5, italic=True, color=GRAY_500)

    add_source(slide, "Fontes: Cano et al. (2017) PMC5360753 · Gachelin et al. (2017) PMC5298425 · Andrade-Neto et al. (2004) PMID14980653 · MDPI Molecules (2020) PMC7571111")


# ============================================================
# Slide 08 — UNICEF Alignment
# ============================================================
def slide08_alignment(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "PARCERIA", "Alinhamento com Prioridades UNICEF", bg_color1=NAVY)

    alignments = [
        ('Proteção infantil', '75% das vítimas mortais africanas são crianças < 5 anos (WHO 2025)',
         'Perfil de segurança pediátrica a estabelecer — incluso no protocolo fase I/II'),
        ('Redução da mortalidade', 'Reduzir mortalidade infantil na África Subsaariana — ODS 3',
         'Se validado clinicamente: potencial de escala às populações mais afectadas'),
        ('Equidade no acesso', '43% da população sem proteção adequada — foco em comunidades remotas',
         'Vantagem logística a confirmar: formulação aquosa sem cadeia de frio'),
        ('Resposta à resistência', 'Resistência à artemisina confirmada em 8 países africanos — necessidade urgente de alternativas',
         'Alcalóides quinínicos da Cinchona: mecanismo de acção distinto da artemisina (documentado)'),
    ]

    card_h = Inches(1.15)
    for i, (priority, desc, match) in enumerate(alignments):
        y = Inches(1.45) + i * Inches(1.3)
        add_rect(slide, Inches(0.5), y, Inches(12.3), card_h, GRAY_50)

        # Priority title
        add_textbox(slide, Inches(0.75), y + Inches(0.1), Inches(5.5), Inches(0.35),
                    priority, font_size=11, bold=True, color=DARK)
        add_textbox(slide, Inches(0.75), y + Inches(0.5), Inches(5.5), Inches(0.5),
                    desc, font_size=9, color=GRAY_500)

        # Match box
        add_rect(slide, Inches(6.5), y + Inches(0.12), Inches(6.1), Inches(0.9), BLUE_50)
        add_textbox(slide, Inches(6.7), y + Inches(0.2), Inches(5.8), Inches(0.65),
                    f"✓ {match}", font_size=9, bold=True, color=NAVY)


# ============================================================
# Slide 09 — Strategy / Pilot Countries
# ============================================================
def slide09_strategy(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "INVESTIGAÇÃO", "Países Piloto para Investigação Clínica",
                           bg_color1=UNICEF_BLUE,
                           subtitle="3 países seleccionados por carga de malária, capacidade regulatória e presença UNICEF")

    countries = [
        ('Nigéria', UNICEF_BLUE,
         "31.9% das mortes africanas — maior carga global. UNICEF e Gavi activos desde 2024. NAFDAC com track record de aprovação ágil."),
        ('Uganda', PINK,
         "Maior rollout de vacina anti-malária de África (Abr 2025, 19º país). Malaria Consortium + UNICEF activos em 25 distritos. Forte capacidade regulatória."),
        ('Moçambique', NAVY,
         "Vacina introduzida em 2024. Sul de África — diversidade geográfica. Plataforma digital upSCALE para suporte a trabalhadores de saúde."),
    ]

    card_h = Inches(1.3)
    for i, (name, color, why) in enumerate(countries):
        y = Inches(1.5) + i * Inches(1.5)
        # Left border
        add_rect(slide, Inches(0.5), y, Inches(0.1), card_h, color)
        # Card background
        add_rect(slide, Inches(0.6), y, Inches(12.1), card_h, BLUE_50)
        # Country name
        add_textbox(slide, Inches(0.9), y + Inches(0.12), Inches(11.5), Inches(0.35),
                    name, font_size=13, bold=True, color=DARK)
        add_textbox(slide, Inches(0.9), y + Inches(0.52), Inches(11.5), Inches(0.65),
                    why, font_size=9, color=GRAY_600)

    # Bottom note
    note_y = Inches(6.1)
    add_rect(slide, Inches(0.5), note_y, Inches(12.3), Inches(0.55), GRAY_50)
    add_textbox(slide, Inches(0.8), note_y + Inches(0.1), Inches(12.0), Inches(0.35),
                "Centros clínicos parceiros em cada país  •  Comités de ética locais e supervisão OMS/UNICEF  •  Plataforma de dados para recolha e análise clínica",
                font_size=9, bold=True, color=DARK)


# ============================================================
# Slide 10 — Timeline
# ============================================================
def slide10_timeline(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "CALENDÁRIO", "Roadmap de Investigação Clínica",
                           bg_color1=NAVY, bg_color2=UNICEF_BLUE,
                           subtitle="Início: Setembro 2026 → Decisão de escala: 2030 (condicionada a resultados)")

    phases = [
        ('SET–DEZ 2026', 'Protocolo & Ética', UNICEF_BLUE, '01',
         ['Desenho do protocolo de investigação fase I/II',
          'Submissão a comités de ética (NG, UG, MZ)',
          'Parceria com universidades e centros clínicos locais']),
        ('JAN–JUN 2027', 'Ensaio Fase I — Segurança', PINK, '02',
         ['Recrutamento de participantes adultos (coorte inicial)',
          'Avaliação de segurança, tolerância e farmacocinética',
          'Supervisão independente OMS / IRB']),
        ('JUL–DEZ 2027', 'Análise Intermédia', NAVY, '03',
         ['Relatório de segurança independente',
          'Decisão baseada em dados: avançar, ajustar ou parar',
          'Partilha de resultados com UNICEF e parceiros']),
        ('2028', 'Ensaio Fase II — Eficácia', UNICEF_BLUE, '04',
         ['Se fase I aprovada: ensaio de eficácia controlado',
          'Expansão a populações vulneráveis (sujeito a dados de segurança)',
          'Primeira publicação peer-review']),
        ('2029', 'Avaliação Regulatória', PINK, '05',
         ['Submissão a SRA (EMA/FDA) ou equivalente regional',
          'Processo WHO Prequalification (se resultados positivos)',
          'Publicações científicas adicionais']),
        ('2030', 'Decisão de Escala', NAVY, '06',
         ['Avaliação de viabilidade para procurement UNICEF',
          'Plano de acesso a países de alta carga (condicionado a aprovação)',
          'Relatório ODS e impacto documentado']),
    ]

    card_w = Inches(4.0)
    card_h = Inches(2.2)

    for i, (period, title, color, phase_num, items) in enumerate(phases):
        row = i // 3
        col = i % 3
        x = Inches(0.5) + col * Inches(4.3)
        y = Inches(1.5) + row * Inches(2.5)

        add_rect(slide, x, y, card_w, card_h, GRAY_50)

        # Phase circle
        add_rect(slide, x - Inches(0.1), y - Inches(0.12), Inches(0.4), Inches(0.4), color)
        add_textbox(slide, x - Inches(0.1), y - Inches(0.12), Inches(0.4), Inches(0.4),
                    phase_num, font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Period
        add_textbox(slide, x + Inches(0.15), y + Inches(0.08), card_w - Inches(0.25), Inches(0.28),
                    period, font_size=8, bold=True, color=color)
        # Title
        add_textbox(slide, x + Inches(0.15), y + Inches(0.38), card_w - Inches(0.25), Inches(0.35),
                    title, font_size=10, bold=True, color=DARK)
        # Items
        items_text = "\n".join([f"› {item}" for item in items])
        add_textbox(slide, x + Inches(0.15), y + Inches(0.78), card_w - Inches(0.25), Inches(1.25),
                    items_text, font_size=8, color=GRAY_600)


# ============================================================
# Slide 11 — Impact / Milestones
# ============================================================
def slide11_impact(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    # Gradient header pink to navy
    header_rect = add_rect(slide, 0, 0, SLIDE_W, Inches(1.3), PINK)
    header_rect.fill.gradient()
    header_rect.fill.gradient_angle = 0
    header_rect.fill.gradient_stops[0].position = 0.0
    header_rect.fill.gradient_stops[0].color.rgb = PINK
    header_rect.fill.gradient_stops[1].position = 1.0
    header_rect.fill.gradient_stops[1].color.rgb = NAVY
    header_rect.line.fill.background()

    add_textbox(slide, Inches(0.8), Inches(0.1), Inches(11), Inches(0.3),
                "RESULTADOS", font_size=8, bold=True, color=RGBColor(0xFF, 0xCC, 0xE0))
    add_textbox(slide, Inches(0.8), Inches(0.38), Inches(11), Inches(0.6),
                "Marcos Científicos e Regulatórios", font_size=22, bold=True, color=WHITE)

    metrics = [
        ('Fase I', 'Dados de segurança e tolerância', '2027 — adultos, coorte inicial', UNICEF_BLUE, BLUE_50),
        ('Fase II', 'Dados de eficácia controlada', '2028 — sujeito a fase I aprovada', PINK, PINK_50),
        ('Reg.', 'Avaliação regulatória WHO', '2029–2030 — se resultados positivos', NAVY, RGBColor(0xEE, 0xF2, 0xFF)),
    ]

    card_w = Inches(3.8)
    card_h = Inches(2.3)

    for i, (value, label, sub, color, bg) in enumerate(metrics):
        x = Inches(0.5) + i * Inches(4.3)
        y = Inches(1.5)
        add_rect(slide, x, y, card_w, card_h, bg)
        add_textbox(slide, x, y + Inches(0.25), card_w, Inches(0.9),
                    value, font_size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.15), card_w, Inches(0.45),
                    label, font_size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.65), card_w, Inches(0.45),
                    sub, font_size=9, color=GRAY_500, align=PP_ALIGN.CENTER)

    # Progress bar visual
    bar_y = Inches(4.1)
    add_rect(slide, Inches(0.5), bar_y, Inches(12.3), Inches(2.0), GRAY_50)
    add_textbox(slide, Inches(0.8), bar_y + Inches(0.12), Inches(12), Inches(0.28),
                "FASES DE INVESTIGAÇÃO — DECISÕES CONDICIONADAS A DADOS", font_size=7, bold=True, color=GRAY_500)

    bars = [
        (5, NAVY),
        (20, NAVY),
        (50, PINK),
        (100, UNICEF_BLUE),
    ]
    years = ['2026', '2027', '2028–29', '2030']
    bar_section_w = Inches(2.8)

    for i, ((h_pct, color), year) in enumerate(zip(bars, years)):
        bx = Inches(0.8) + i * Inches(3.0)
        max_bar_h = Inches(0.85)
        actual_bar_h = max_bar_h * h_pct / 100
        bar_bottom = bar_y + Inches(1.6)
        add_rect(slide, bx, bar_bottom - actual_bar_h, bar_section_w - Inches(0.3), actual_bar_h, color)
        add_textbox(slide, bx, bar_y + Inches(1.65), bar_section_w - Inches(0.3), Inches(0.28),
                    year, font_size=8, color=GRAY_400, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 12 — Partnership (full-color)
# ============================================================
def slide12_partnership(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_gradient_background(slide, UNICEF_BLUE, NAVY)

    # Header text
    add_textbox(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(0.3),
                "PROPOSTA", font_size=8, color=BLUE_200, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(0.85), Inches(11.3), Inches(0.75),
                "All Shanti + UNICEF", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.45),
                "Uma parceria de investigação para determinar o potencial clínico do All Shanti",
                font_size=12, color=BLUE_200, align=PP_ALIGN.CENTER)

    # Pink line
    add_rect(slide, Inches(6.1), Inches(2.1), Inches(1.1), Inches(0.07), PINK)

    pillars = [
        ('Co-investigação clínica', 'Protocolo conjunto fase I/II com supervisão UNICEF/OMS — rigor científico e independência garantidos'),
        ('Co-financiamento de ensaios', 'Acesso a fundos de investigação (Gavi, GFATM, PMI) para conduzir os ensaios clínicos necessários'),
        ('Recolha e análise de dados', 'Sistemas UNICEF de monitorização para recolha robusta de dados clínicos em contexto africano'),
        ('Credibilidade regulatória', 'O envolvimento UNICEF reforça a credibilidade junto de SRAs e acelera o processo WHO Prequalification'),
    ]

    card_w = Inches(5.8)
    card_h = Inches(1.2)

    for i, (title, desc) in enumerate(pillars):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.5)
        y = Inches(2.4) + row * Inches(1.4)

        box = add_rect(slide, x, y, card_w, card_h, NAVY)
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x0A, 0x3A, 0x8F)
        box.line.fill.background()

        add_textbox(slide, x + Inches(0.2), y + Inches(0.12), card_w - Inches(0.3), Inches(0.35),
                    title, font_size=11, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.5), card_w - Inches(0.3), Inches(0.55),
                    desc, font_size=8.5, color=BLUE_200)

    # Objective
    add_textbox(slide, Inches(0.8), Inches(5.45), Inches(11.7), Inches(0.7),
                "Objectivo: Determinar, com rigor científico, o potencial do All Shanti como ferramenta antimalárica — e, se validado, levá-lo a quem mais precisa",
                font_size=10, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 13 — Next Steps
# ============================================================
def slide13_nextsteps(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background_fill(slide, WHITE)

    header_h = add_header(slide, "ACÇÃO", "Próximos Passos", bg_color1=NAVY)

    steps = [
        ('01', 'Reunião técnica',
         'Apresentação da revisão bibliográfica e plano de investigação clínica ao departamento de saúde da UNICEF',
         'Abril 2026', UNICEF_BLUE),
        ('02', 'Avaliação de parceria',
         'Avaliação formal do potencial de parceria UNICEF–All Shanti e definição de termos',
         'Maio 2026', PINK),
        ('03', 'Países piloto prioritários',
         'Definição final dos 3 países piloto e início das diligências regulatórias locais',
         'Jun–Jul 2026', NAVY),
        ('04', 'MOU & Financiamento',
         'Assinatura de Memorando de Entendimento e identificação de fontes de financiamento (Gavi, GFATM)',
         'Ago–Set 2026', UNICEF_BLUE),
    ]

    card_h = Inches(1.1)
    for i, (num, title, desc, deadline, color) in enumerate(steps):
        y = Inches(1.45) + i * Inches(1.35)

        # Number circle
        circle = add_rect(slide, Inches(0.5), y + Inches(0.15), Inches(0.65), Inches(0.65), color)
        add_textbox(slide, Inches(0.5), y + Inches(0.15), Inches(0.65), Inches(0.65),
                    num, font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Card
        add_rect(slide, Inches(1.35), y, Inches(11.2), card_h, GRAY_50)
        add_textbox(slide, Inches(1.6), y + Inches(0.1), Inches(8.5), Inches(0.35),
                    title, font_size=11, bold=True, color=DARK)
        add_textbox(slide, Inches(1.6), y + Inches(0.5), Inches(8.5), Inches(0.5),
                    desc, font_size=9, color=GRAY_500)

        # Deadline badge
        add_rect(slide, Inches(10.0), y + Inches(0.25), Inches(2.3), Inches(0.4), color)
        add_textbox(slide, Inches(10.0), y + Inches(0.25), Inches(2.3), Inches(0.4),
                    deadline, font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ============================================================
# Slide 14 — Thanks (full-color)
# ============================================================
def slide14_thanks(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    add_gradient_background(slide, NAVY, UNICEF_BLUE)

    # AS circle
    logo = add_rect(slide, Inches(6.1), Inches(1.0), Inches(1.1), Inches(1.1), PINK)
    add_textbox(slide, Inches(6.1), Inches(1.0), Inches(1.1), Inches(1.1),
                "AS", font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Thank you
    add_textbox(slide, Inches(1), Inches(2.25), Inches(11.3), Inches(1.0),
                "Obrigado", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.5),
                "All Shanti — Em investigação para combater a malária", font_size=14, color=BLUE_200, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.4),
                "Águas de São Silvestre, SA", font_size=12, color=BLUE_200, align=PP_ALIGN.CENTER)

    # Pink line
    add_rect(slide, Inches(6.1), Inches(4.3), Inches(1.1), Inches(0.07), PINK)

    # Contact info
    add_textbox(slide, Inches(1), Inches(4.55), Inches(11.3), Inches(0.35),
                "Águas de São Silvestre, SA", font_size=10, color=BLUE_200, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.9), Inches(11.3), Inches(0.35),
                "Pernes, Santarém — Portugal", font_size=9, color=BLUE_200, align=PP_ALIGN.CENTER)

    # Footer
    add_textbox(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.35),
                "APRESENTAÇÃO CONFIDENCIAL — UNICEF 2026", font_size=8, color=BLUE_200, align=PP_ALIGN.CENTER)


# ============================================================
# MAIN
# ============================================================
def main():
    import pptx.enum.shapes

    prs = create_presentation()

    print("Creating slides...")
    slide01_cover(prs)
    print("  Slide 01 Cover — done")

    slide02_crisis(prs)
    print("  Slide 02 Crisis — done")

    slide03_countries(prs)
    print("  Slide 03 Countries — done")

    slide04_economic(prs)
    print("  Slide 04 Economic — done")

    slide05_treatment(prs)
    print("  Slide 05 Treatment — done")

    slide06_solution(prs)
    print("  Slide 06 Solution — done")

    slide07_advantages(prs)
    print("  Slide 07 Advantages — done")

    slide_science(prs)
    print("  Slide Science — done")

    slide08_alignment(prs)
    print("  Slide 08 Alignment — done")

    slide09_strategy(prs)
    print("  Slide 09 Strategy — done")

    slide10_timeline(prs)
    print("  Slide 10 Timeline — done")

    slide11_impact(prs)
    print("  Slide 11 Impact — done")

    slide12_partnership(prs)
    print("  Slide 12 Partnership — done")

    slide13_nextsteps(prs)
    print("  Slide 13 Next Steps — done")

    slide14_thanks(prs)
    print("  Slide 14 Thanks — done")

    output_path = "/Users/gustavoteixeiracuco/Documents/github/all-shanti-unicef-pitch/All_Shanti_UNICEF_Pitch_2026.pptx"
    prs.save(output_path)
    print(f"\nSaved: {output_path}")

    import os
    size = os.path.getsize(output_path)
    print(f"File size: {size:,} bytes ({size/1024:.1f} KB)")


if __name__ == "__main__":
    main()
